from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors (RGB)
BG_COLOR = RGBColor(0x1A, 0x0B, 0x2E) # Deep Purple
ACCENT_COLOR = RGBColor(0xF5, 0x9E, 0x0B) # Gold/Amber
TEXT_MAIN = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SUB = RGBColor(0xE9, 0xD5, 0xFF) # Light Purple

def create_deck():
    prs = Presentation()
    
    # helper to set background
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "SajuChain"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Destiny on Chain: Where Ancient Wisdom Meets Eternal Ledger\nInvestor Presentation 2026"
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # Helper for Content Slides
    def add_content_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_bg(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = content_text
        
        for p in tf.paragraphs:
            p.font.color.rgb = TEXT_MAIN
            p.font.size = Pt(18)
            
    # 2. Problem
    add_content_slide(
        "The Problem: Uncertainty & Ephemerality", 
        "• Subjective Interpretation: Relies on individual skill, inconsistent results.\n"
        "• Lack of Permanence: Physical amulets lost. No digital record.\n"
        "• Data Privacy: Sensitive birth data shared with unverified services."
    )

    # 3. Solution
    add_content_slide(
        "The Solution: SajuChain",
        "1. Algorithmic Precision: lunar-javascript + GPT-4o.\n"
        "2. Immutable Ownership: Minted as Compressed NFTs (cNFT) on Solana.\n"
        "3. Digital Mysticism: Ritualistic UI/UX."
    )

    # 4. Tech
    add_content_slide(
        "Technical Architecture",
        "• Frontend: Next.js 16, Tailwind CSS v4, Framer Motion\n"
        "• Blockchain: Solana Web3.js, Metaplex (Umi)\n"
        "• AI Engine: Hybrid Logic (Deterministic + Generative)"
    )

    # 5. Roadmap
    add_content_slide(
        "Roadmap & Vision",
        "• Phase 1 (Now): MVP, Devnet, UI/UX V2 (Completed)\n"
        "• Phase 2 (Q3 2026): Mainnet, Mobile App\n"
        "• Phase 3 (2027): Oracle DAO, Metaverse Identity"
    )

    prs.save('sajuchain_pitch_deck.pptx')
    print("Created sajuchain_pitch_deck.pptx")

if __name__ == "__main__":
    create_deck()
